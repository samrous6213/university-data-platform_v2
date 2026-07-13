import os
import sys
import hashlib
import logging
from datetime import datetime
from io import BytesIO

import fitz
import pandas as pd

from docx import Document
from pptx import Presentation

from pyspark.sql import SparkSession
from pyspark.sql.types import *

PROJECT_ROOT = os.getenv(
    "PROJECT_ROOT",
    "/opt/spark/work-dir"
)

if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

from src.storage.minio.hiba_client import MinIOClient

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

logger = logging.getLogger(__name__)

BUCKET = "raw-documents"

OUTPUT_PATH = os.getenv(
    "DOCUMENT_PARSER_OUTPUT",
    "/opt/spark/work-dir/data/processed_documents"
)

SPARK_MASTER = os.getenv(
    "SPARK_MASTER",
    "spark://spark-master:7077"
)

CRAWLER_VERSION = os.getenv(
    "CRAWLER_VERSION",
    "chaimae_pdf_parser_v1"
)

MAX_TEXT_LEN = 100000

SCHEMA = StructType([
    StructField("record_id", StringType(), False),
    StructField("source_system", StringType(), True),
    StructField("storage_path", StringType(), False),
    StructField("file_name", StringType(), True),
    StructField("document_type", StringType(), True),
    StructField("extracted_text", StringType(), True),
    StructField("content_hash", StringType(), True),
    StructField("crawler_version", StringType(), True),
    StructField("crawl_timestamp", StringType(), False),
])


def extract_source_system(path):

    for part in path.split("/"):

        if part.startswith("source="):
            return part.replace("source=", "")

    return "unknown"


def compute_content_hash(file_bytes):

    return hashlib.md5(
        file_bytes
    ).hexdigest()


def extract_pdf_text(file_bytes):

    text = ""

    try:

        pdf = fitz.open(
            stream=file_bytes,
            filetype="pdf"
        )

        for page in pdf:
            text += page.get_text() + "\n"

        pdf.close()

    except Exception as e:

        logger.error(
            f"PDF parse error: {e}"
        )

    return text[:MAX_TEXT_LEN]


def extract_docx_text(file_bytes):

    text = ""

    try:

        doc = Document(
            BytesIO(file_bytes)
        )

        for paragraph in doc.paragraphs:

            text += (
                paragraph.text + "\n"
            )

    except Exception as e:

        logger.error(
            f"DOCX parse error: {e}"
        )

    return text[:MAX_TEXT_LEN]


def extract_excel_text(file_bytes):

    text = ""

    try:

        excel = pd.ExcelFile(
            BytesIO(file_bytes)
        )

        for sheet in excel.sheet_names:

            df = pd.read_excel(
                BytesIO(file_bytes),
                sheet_name=sheet
            )

            text += (
                f"\n=== {sheet} ===\n"
            )

            text += (
                df.astype(str)
                .to_string()
            )

            text += "\n"

    except Exception as e:

        logger.error(
            f"Excel parse error: {e}"
        )

    return text[:MAX_TEXT_LEN]


def extract_csv_text(file_bytes):

    text = ""

    try:

        df = pd.read_csv(
            BytesIO(file_bytes)
        )

        text += (
            df.astype(str)
            .to_string()
        )

    except Exception as e:

        logger.error(
            f"CSV parse error: {e}"
        )

    return text[:MAX_TEXT_LEN]


def extract_pptx_text(file_bytes):

    text = ""

    try:

        presentation = Presentation(
            BytesIO(file_bytes)
        )

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += (
                        shape.text + "\n"
                    )

    except Exception as e:

        logger.error(
            f"PPTX parse error: {e}"
        )

    return text[:MAX_TEXT_LEN]


def build_record(
    object_name,
    text,
    document_type,
    content_hash
):

    return {

        "record_id":
            hashlib.md5(
                object_name.encode(
                    "utf-8"
                )
            ).hexdigest(),

        "source_system":
            extract_source_system(
                object_name
            ),

        "storage_path":
            object_name,

        "file_name":
            os.path.basename(
                object_name
            ),

        "document_type":
            document_type,

        "extracted_text":
            text,

        "content_hash":
            content_hash,

        "crawler_version":
            CRAWLER_VERSION,

        "crawl_timestamp":
            datetime.now().isoformat()
    }


def main():

    spark = (
        SparkSession.builder
        .master(SPARK_MASTER)
        .appName(
            "Universal Document Parser"
        )
        .getOrCreate()
    )

    logger.info(
        "Connecting to MinIO..."
    )

    client = MinIOClient(
        endpoint=os.getenv(
            "MINIO_ENDPOINT",
            "university-minio:9000"
        ),
        access_key=os.getenv(
            "MINIO_ACCESS_KEY",
            "minioadmin"
        ),
        secret_key=os.getenv(
            "MINIO_SECRET_KEY",
            "minioadmin"
        ),
        secure=os.getenv(
            "MINIO_SECURE",
            "false"
        ).lower() == "true"
    )

    records = []

    processed = 0

    logger.info(
        f"Scanning bucket: {BUCKET}"
    )

    for obj in client.client.list_objects(
        BUCKET,
        recursive=True
    ):

        extension = (
            os.path.splitext(
                obj.object_name
            )[1]
            .lower()
        )

        if extension not in [
            ".pdf",
            ".docx",
            ".xlsx",
            ".xls",
            ".pptx",
            ".csv"
        ]:
            logger.info(
                f"Skipping unsupported file: {obj.object_name}"
            )
            continue

        logger.info(
            f"Processing document: {obj.object_name}"
        )

        try:

            response = client.client.get_object(
                BUCKET,
                obj.object_name
            )

            file_bytes = response.read()

            response.close()
            response.release_conn()

            content_hash = compute_content_hash(
                file_bytes
            )

            text = ""

            if extension == ".pdf":

                text = extract_pdf_text(
                    file_bytes
                )

            elif extension == ".docx":

                text = extract_docx_text(
                    file_bytes
                )

            elif extension in [
                ".xlsx",
                ".xls"
            ]:

                text = extract_excel_text(
                    file_bytes
                )

            elif extension == ".pptx":

                text = extract_pptx_text(
                    file_bytes
                )

            elif extension == ".csv":

                text = extract_csv_text(
                    file_bytes
                )

            records.append(
                build_record(
                    obj.object_name,
                    text,
                    extension.replace(
                        ".",
                        ""
                    ),
                    content_hash
                )
            )

            processed += 1

            if processed % 25 == 0:

                logger.info(
                    f"{processed} documents processed..."
                )

        except Exception as e:

            logger.error(
                f"Failed: {obj.object_name}"
            )

            logger.error(str(e))

    logger.info(
        f"Total documents = {processed}"
    )

    if not records:

        logger.warning(
            "No documents found. Stopping Spark gracefully."
        )

        spark.stop()
        return

    df = spark.createDataFrame(
        records,
        schema=SCHEMA
    )

    logger.info(
        "Spark schema:"
    )

    df.printSchema()

    logger.info(
        "Sample rows:"
    )

    df.show(
        10,
        truncate=False
    )

    (
        df.write
        .mode("overwrite")
        .parquet(
            OUTPUT_PATH
        )
    )

    logger.info(
        f"Saved to {OUTPUT_PATH}"
    )

    logger.info(
        f"Rows = {df.count()}"
    )

    logger.info(
        "Document parsing completed successfully."
    )

    spark.stop()


if __name__ == "__main__":
    main()